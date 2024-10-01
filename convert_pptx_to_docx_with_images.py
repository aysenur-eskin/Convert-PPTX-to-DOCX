import os
from tkinter import Tk, Button, Label, filedialog
from pptx import Presentation
from docx import Document
import pytesseract
from PIL import Image
import tempfile
import re

def clean_text(text):
    cleaned_text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', text)
    return cleaned_text

def ocr_image_to_text(image_path):
    image = Image.open(image_path)
    text = pytesseract.image_to_string(image, lang="tur")
    return text

def extract_text_from_images(pptx_path):
    extracted_texts = []
    presentation = Presentation(pptx_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                temp_image = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                temp_image.write(image_bytes)
                temp_image.close()
                text = ocr_image_to_text(temp_image.name)
                extracted_texts.append(clean_text(text))                
                os.remove(temp_image.name)
    return extracted_texts

def convert_pptx_to_docx(pptx_path, docx_path):
    presentation = Presentation(pptx_path)
    document = Document()
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                cleaned_text = clean_text(shape.text)
                document.add_paragraph(cleaned_text)
            elif hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                temp_image = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                temp_image.write(image_bytes)
                temp_image.close()
                text = ocr_image_to_text(temp_image.name)
                document.add_paragraph(clean_text(text))
                os.remove(temp_image.name)
    document.save(docx_path)

def select_input_files():
    file_paths = filedialog.askopenfilenames(filetypes=[("PowerPoint files", "*.pptx")])
    input_files_label.config(text=", ".join(file_paths))
    return file_paths

def convert_to_docx():
    file_paths = input_files_label.cget("text").split(", ")
    for file_path in file_paths:
        if file_path.endswith('.pptx'):
            output_path = filedialog.asksaveasfilename(defaultextension=".docx", 
            filetypes=[("Word files", "*.docx")])
            if output_path:
                convert_pptx_to_docx(file_path, output_path)
                print(f'{os.path.basename(file_path)} converted to {os.path.basename(output_path)} successfully.')
    status_label.config(text="Conversion complete.")

root = Tk()
root.title("PPTX to DOCX Converter")
input_files_label = Label(root, text="Select input files")
input_files_label.grid(row=0, column=0)
status_label = Label(root, text="")
status_label.grid(row=2, column=0)
input_files_button = Button(root, text="Browse", command=select_input_files)
input_files_button.grid(row=0, column=1)
convert_button = Button(root, text="Convert", command=convert_to_docx)
convert_button.grid(row=1, column=0, columnspan=2)
root.mainloop()