import os
from tkinter import Tk, Button, Label, filedialog
from pptx import Presentation
from docx import Document
import re

def clean_text(text):
    cleaned_text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', text)
    return cleaned_text

def pptx_to_docx(file_paths):
    for file_path in file_paths:
        if file_path.endswith('.pptx'):
            output_path = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("Word files", "*.docx")])
            if output_path:
                convert_pptx_to_docx(file_path, output_path)
                print(f'{os.path.basename(file_path)} converted to {os.path.basename(output_path)} successfully.')

def convert_pptx_to_docx(pptx_path, docx_path):
    presentation = Presentation(pptx_path)
    document = Document()

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                cleaned_text = clean_text(shape.text)
                document.add_paragraph(cleaned_text)

    document.save(docx_path)

def select_input_files():
    file_paths = filedialog.askopenfilenames(filetypes=[("PowerPoint files", "*.pptx")])
    input_files_label.config(text=", ".join(file_paths))
    return file_paths

def convert_to_docx():
    file_paths = input_files_label.cget("text").split(", ")
    pptx_to_docx(file_paths)
    status_label.config(text="Conversion complete.")

root = Tk()
root.title("PPTX to DOCX Converter")
input_files_label = Label(root, text="Select input files")
input_files_label.grid(row=0, column=0)
status_label = Label(root, text="")
status_label.grid(row=2, column=0)
input_files_button = Button(root, text="Browse", command=select_input_files)
input_files_button.grid(row=0, column=1)
convert_button = Button(root, text="Convert", command=convert_to_docx)
convert_button.grid(row=1, column=0, columnspan=2)
root.mainloop()
